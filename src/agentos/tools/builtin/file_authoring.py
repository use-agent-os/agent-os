"""Channel-safe structured file authoring tools."""

from __future__ import annotations

import csv
import hashlib
import io
import json
import zipfile
from pathlib import Path
from typing import Any
from xml.sax.saxutils import escape

from agentos.artifacts import (
    DEFAULT_ARTIFACT_DISK_BUDGET_BYTES,
    DEFAULT_ARTIFACT_MAX_BYTES,
    ArtifactBudgetError,
    ArtifactStore,
    artifact_payload,
)
from agentos.tools.registry import tool
from agentos.tools.types import ToolError, current_tool_context

_CSV_MIME = "text/csv"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_PDF_MIME = "application/pdf"
_PDF_SANS_FONT = "AgentOSPDFSans"
_PDF_SANS_BOLD_FONT = "AgentOSPDFSans-Bold"
_PDF_CJK_FONT = "STSong-Light"
_PDF_SANS_CANDIDATES = (
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/local/share/fonts/dejavu/DejaVuSans.ttf",
    "/Library/Fonts/Arial Unicode.ttf",
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    "C:/Windows/Fonts/arial.ttf",
)
_STABLE_ZIP_TIMESTAMP = (1980, 1, 1, 0, 0, 0)
_PDF_SANS_BOLD_CANDIDATES = (
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/local/share/fonts/dejavu/DejaVuSans-Bold.ttf",
    "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
    "C:/Windows/Fonts/arialbd.ttf",
)
_PDF_FONT_REGISTERED = False


def _ensure_name(name: str | None, *, default: str, suffix: str) -> str:
    raw = Path(str(name or default)).name.strip()
    if not raw or raw in {".", ".."}:
        raw = default
    if not raw.lower().endswith(suffix):
        raw = f"{raw}{suffix}"
    return raw


def _stringify_cell(value: Any) -> str | int | float | bool | None:
    if value is None or isinstance(value, str | int | float | bool):
        return value
    return json.dumps(value, ensure_ascii=False, sort_keys=True)


def _rows(value: Any, *, field_name: str) -> list[list[str | int | float | bool | None]]:
    if not isinstance(value, list) or not value:
        raise ToolError(f"{field_name} must be a non-empty list of rows")
    rows: list[list[str | int | float | bool | None]] = []
    for index, row in enumerate(value, start=1):
        if not isinstance(row, list):
            raise ToolError(f"{field_name}[{index}] must be a list")
        rows.append([_stringify_cell(cell) for cell in row])
    return rows


def _sheet_title(value: Any, fallback: str) -> str:
    raw = str(value or fallback).strip()
    cleaned = "".join("_" if char in "[]:*?/\\" else char for char in raw)
    cleaned = cleaned.strip("'") or fallback
    return cleaned[:31]


def _text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    if isinstance(value, list):
        return "\n".join(_text(item) for item in value)
    if isinstance(value, dict):
        return json.dumps(value, ensure_ascii=False, sort_keys=True)
    return str(value)


def _first_existing_path(candidates: tuple[str, ...]) -> str | None:
    for candidate in candidates:
        if Path(candidate).exists():
            return candidate
    return None


def _normalize_zip_timestamps(payload: bytes) -> bytes:
    source = io.BytesIO(payload)
    target = io.BytesIO()
    with zipfile.ZipFile(source, "r") as src, zipfile.ZipFile(target, "w") as dst:
        for info in src.infolist():
            stable = zipfile.ZipInfo(info.filename, _STABLE_ZIP_TIMESTAMP)
            stable.compress_type = info.compress_type
            stable.external_attr = info.external_attr
            stable.comment = info.comment
            dst.writestr(stable, src.read(info.filename))
    return target.getvalue()


def _register_pdf_fonts() -> tuple[str, str, str | None]:
    global _PDF_FONT_REGISTERED

    from reportlab.pdfbase import pdfmetrics  # type: ignore[import-untyped]
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont  # type: ignore[import-untyped]
    from reportlab.pdfbase.ttfonts import TTFont  # type: ignore[import-untyped]

    if not _PDF_FONT_REGISTERED:
        sans_path = _first_existing_path(_PDF_SANS_CANDIDATES)
        if sans_path is not None and _PDF_SANS_FONT not in pdfmetrics.getRegisteredFontNames():
            pdfmetrics.registerFont(TTFont(_PDF_SANS_FONT, sans_path))

        bold_path = _first_existing_path(_PDF_SANS_BOLD_CANDIDATES)
        if bold_path is not None and _PDF_SANS_BOLD_FONT not in pdfmetrics.getRegisteredFontNames():
            pdfmetrics.registerFont(TTFont(_PDF_SANS_BOLD_FONT, bold_path))

        if _PDF_CJK_FONT not in pdfmetrics.getRegisteredFontNames():
            pdfmetrics.registerFont(UnicodeCIDFont(_PDF_CJK_FONT))

        _PDF_FONT_REGISTERED = True

    font_names = set(pdfmetrics.getRegisteredFontNames())
    base_font = _PDF_SANS_FONT if _PDF_SANS_FONT in font_names else "Helvetica"
    bold_font = _PDF_SANS_BOLD_FONT if _PDF_SANS_BOLD_FONT in font_names else base_font
    cjk_font = _PDF_CJK_FONT if _PDF_CJK_FONT in font_names else None
    return base_font, bold_font, cjk_font


def _is_cjk(char: str) -> bool:
    codepoint = ord(char)
    return (
        0x3400 <= codepoint <= 0x4DBF
        or 0x4E00 <= codepoint <= 0x9FFF
        or 0xF900 <= codepoint <= 0xFAFF
        or 0x20000 <= codepoint <= 0x2A6DF
        or 0x2A700 <= codepoint <= 0x2B73F
        or 0x2B740 <= codepoint <= 0x2B81F
        or 0x2B820 <= codepoint <= 0x2CEAF
        or 0x2CEB0 <= codepoint <= 0x2EBEF
        or 0x30000 <= codepoint <= 0x3134F
    )


def _font_supports_char(font_name: str, char: str) -> bool:
    from reportlab.pdfbase import pdfmetrics  # type: ignore[import-untyped]

    font = pdfmetrics.getFont(font_name)
    face = getattr(font, "face", None)
    char_widths = getattr(face, "charWidths", None)
    if char_widths is None:
        return ord(char) < 256
    return ord(char) in char_widths


def _pdf_markup_text(value: Any, *, base_font: str, cjk_font: str | None) -> str:
    text = _text(value)
    if not text:
        return ""

    parts: list[str] = []
    run: list[str] = []
    run_font: str | None = None

    def flush() -> None:
        nonlocal run, run_font
        if not run:
            return
        escaped = escape("".join(run))
        if run_font is not None:
            parts.append(f'<font name="{run_font}">{escaped}</font>')
        else:
            parts.append(escaped)
        run = []
        run_font = None

    for char in text:
        target_font = cjk_font if cjk_font is not None and _is_cjk(char) else None
        if target_font is None and not _font_supports_char(base_font, char):
            continue
        if target_font != run_font:
            flush()
            run_font = target_font
        run.append(char)
    flush()
    return "".join(parts)


def _published_response(
    *,
    payload: bytes,
    name: str,
    mime: str,
    source: str,
) -> str:
    ctx = current_tool_context.get()
    if ctx is None:
        raise ToolError(f"{source} requires tool context")
    if not ctx.artifact_media_root:
        raise ToolError("artifact storage is not configured for this turn")
    if not ctx.artifact_session_id or not ctx.session_key:
        raise ToolError("artifact session scope is not configured for this turn")

    target_sha256 = hashlib.sha256(payload).hexdigest()
    for published in reversed(ctx.published_artifacts):
        if published.get("sha256") != target_sha256:
            continue
        llm_artifact = {k: v for k, v in published.items() if k != "download_url"}
        return json.dumps(
            {
                "status": "already_published",
                "artifact": llm_artifact,
                "note": (
                    "This generated file is already registered for the current surface "
                    "in this turn. Do not recreate or paste the file contents; just "
                    "confirm it is ready."
                ),
            },
            ensure_ascii=False,
        )

    store = ArtifactStore(ctx.artifact_media_root)
    existing = store.find_existing_ref(
        session_id=ctx.artifact_session_id,
        session_key=ctx.session_key,
        sha256=target_sha256,
        name=name,
        mime=mime,
    )
    if existing is not None:
        artifact = artifact_payload(existing)
        if not any(item.get("id") == artifact.get("id") for item in ctx.published_artifacts):
            ctx.published_artifacts.append(artifact)
        llm_artifact = {k: v for k, v in artifact.items() if k != "download_url"}
        return json.dumps(
            {
                "status": "already_published",
                "artifact": llm_artifact,
                "note": (
                    "This session already has the same generated file registered. "
                    "Do not recreate or republish it; just confirm it is ready."
                ),
            },
            ensure_ascii=False,
        )
    try:
        ref = store.publish_bytes(
            payload,
            session_id=ctx.artifact_session_id,
            session_key=ctx.session_key,
            name=name,
            mime=mime,
            source=source,
            max_bytes=ctx.artifact_max_bytes
            if ctx.artifact_max_bytes is not None
            else DEFAULT_ARTIFACT_MAX_BYTES,
            disk_budget_bytes=ctx.artifact_disk_budget_bytes
            if ctx.artifact_disk_budget_bytes is not None
            else DEFAULT_ARTIFACT_DISK_BUDGET_BYTES,
        )
    except ArtifactBudgetError as exc:
        raise ToolError(str(exc)) from exc

    artifact = artifact_payload(ref)
    ctx.published_artifacts.append(artifact)
    llm_artifact = {k: v for k, v in artifact.items() if k != "download_url"}
    return json.dumps(
        {
            "status": "published",
            "artifact": llm_artifact,
            "note": (
                "The generated file is registered for the active surface. The UI or "
                "channel adapter handles download chips or native delivery; do not "
                "include any URL or paste file source in your reply."
            ),
        },
        ensure_ascii=False,
    )


@tool(
    name="create_csv",
    description=(
        "Create a CSV file from structured rows and publish it as a generated artifact. "
        "Use this for channel file requests instead of writing raw files or pasting CSV text."
    ),
    params={
        "name": {"type": "string", "description": "Output filename. .csv is appended if missing."},
        "rows": {
            "type": "array",
            "description": (
                "Non-empty array of row arrays. Values may be strings, numbers, "
                "booleans, null, arrays, or objects."
            ),
            "items": {"type": "array"},
        },
    },
    required=["rows"],
)
async def create_csv(rows: list[list[Any]], name: str | None = None) -> str:
    output = io.StringIO(newline="")
    writer = csv.writer(output, lineterminator="\n")
    writer.writerows(_rows(rows, field_name="rows"))
    return _published_response(
        payload=output.getvalue().encode("utf-8-sig"),
        name=_ensure_name(name, default="generated.csv", suffix=".csv"),
        mime=_CSV_MIME,
        source="create_csv",
    )


@tool(
    name="create_xlsx",
    description=(
        "Create an XLSX workbook from structured sheets and publish it as a generated artifact. "
        "Use this for spreadsheet requests from channels."
    ),
    params={
        "name": {"type": "string", "description": "Output filename. .xlsx is appended if missing."},
        "sheets": {
            "type": "array",
            "description": "Non-empty array of objects with optional name and required rows.",
            "items": {"type": "object"},
        },
    },
    required=["sheets"],
)
async def create_xlsx(sheets: list[dict[str, Any]], name: str | None = None) -> str:
    if not isinstance(sheets, list) or not sheets:
        raise ToolError("sheets must be a non-empty list")
    try:
        from openpyxl import Workbook  # type: ignore[import-untyped]
    except ImportError as exc:
        raise ToolError("create_xlsx requires openpyxl to be installed") from exc

    workbook = Workbook()
    for index, sheet in enumerate(sheets):
        if not isinstance(sheet, dict):
            raise ToolError(f"sheets[{index + 1}] must be an object")
        worksheet = workbook.active if index == 0 else workbook.create_sheet()
        worksheet.title = _sheet_title(sheet.get("name"), f"Sheet{index + 1}")
        for row in _rows(sheet.get("rows"), field_name=f"sheets[{index + 1}].rows"):
            worksheet.append(row)

    output = io.BytesIO()
    workbook.save(output)
    return _published_response(
        payload=output.getvalue(),
        name=_ensure_name(name, default="generated.xlsx", suffix=".xlsx"),
        mime=_XLSX_MIME,
        source="create_xlsx",
    )


@tool(
    name="create_pptx",
    description=(
        "Create a basic text-only PowerPoint deck from structured slides and publish it as a "
        "generated artifact. This fallback only supports slide titles plus body text or bullets; "
        "it does not support images, charts, icons, custom layouts, or templates."
    ),
    params={
        "name": {"type": "string", "description": "Output filename. .pptx is appended if missing."},
        "slides": {
            "type": "array",
            "description": (
                "Non-empty array of slide objects with title plus body or bullets only. "
                "Images, charts, icons, and custom layout fields are ignored because this "
                "tool is a text-only fallback."
            ),
            "items": {"type": "object"},
        },
    },
    required=["slides"],
    exposed_by_default=False,
)
async def create_pptx(slides: list[dict[str, Any]], name: str | None = None) -> str:
    if not isinstance(slides, list) or not slides:
        raise ToolError("slides must be a non-empty list")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise ToolError("create_pptx requires python-pptx to be installed") from exc

    presentation = Presentation()
    for index, slide_payload in enumerate(slides):
        if not isinstance(slide_payload, dict):
            raise ToolError(f"slides[{index + 1}] must be an object")
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title = _text(slide_payload.get("title")) or f"Slide {index + 1}"
        slide.shapes.title.text = title

        body_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        bullets = slide_payload.get("bullets")
        body = slide_payload.get("body")
        if isinstance(bullets, list):
            lines = [_text(item) for item in bullets if _text(item)]
        else:
            lines = [line for line in _text(body).splitlines() if line.strip()]
        if not lines:
            lines = [""]

        if body_shape is None:
            body_shape = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
        text_frame = body_shape.text_frame
        text_frame.clear()
        for line_index, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if line_index == 0 else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.font.size = Pt(20)

    output = io.BytesIO()
    presentation.save(output)
    return _published_response(
        payload=_normalize_zip_timestamps(output.getvalue()),
        name=_ensure_name(name, default="generated.pptx", suffix=".pptx"),
        mime=_PPTX_MIME,
        source="create_pptx",
    )


@tool(
    name="create_pdf_report",
    description=(
        "Create a simple PDF report from structured text sections and publish it as a "
        "generated artifact. "
        "Use this for channel PDF requests instead of returning PDF source text."
    ),
    params={
        "name": {"type": "string", "description": "Output filename. .pdf is appended if missing."},
        "title": {"type": "string", "description": "Report title."},
        "sections": {
            "type": "array",
            "description": "Optional array of objects with heading and body fields.",
            "items": {"type": "object"},
        },
        "body": {"type": "string", "description": "Optional fallback body text."},
    },
    required=["title"],
)
async def create_pdf_report(
    title: str,
    sections: list[dict[str, Any]] | None = None,
    body: str | None = None,
    name: str | None = None,
) -> str:
    try:
        from reportlab.lib.pagesizes import letter  # type: ignore[import-untyped]
        from reportlab.lib.styles import getSampleStyleSheet  # type: ignore[import-untyped]
        from reportlab.platypus import (  # type: ignore[import-untyped]
            Paragraph,
            SimpleDocTemplate,
            Spacer,
        )
    except ImportError as exc:
        raise ToolError("create_pdf_report requires reportlab to be installed") from exc

    output = io.BytesIO()
    doc = SimpleDocTemplate(output, pagesize=letter)
    styles = getSampleStyleSheet()
    base_font, bold_font, cjk_font = _register_pdf_fonts()
    for style_name in ("Title", "Heading2", "BodyText"):
        styles[style_name].fontName = (
            bold_font if style_name in {"Title", "Heading2"} else base_font
        )

    story: list[Any] = [
        Paragraph(
            _pdf_markup_text(title, base_font=base_font, cjk_font=cjk_font) or "Report",
            styles["Title"],
        ),
        Spacer(1, 12),
    ]

    if sections:
        if not isinstance(sections, list):
            raise ToolError("sections must be a list when provided")
        for index, section in enumerate(sections):
            if not isinstance(section, dict):
                raise ToolError(f"sections[{index + 1}] must be an object")
            heading = _text(section.get("heading")) or f"Section {index + 1}"
            section_body = _text(section.get("body"))
            story.append(
                Paragraph(
                    _pdf_markup_text(heading, base_font=base_font, cjk_font=cjk_font),
                    styles["Heading2"],
                )
            )
            if section_body:
                for paragraph in section_body.splitlines():
                    if paragraph.strip():
                        story.append(
                            Paragraph(
                                _pdf_markup_text(paragraph, base_font=base_font, cjk_font=cjk_font),
                                styles["BodyText"],
                            )
                        )
            story.append(Spacer(1, 8))
    elif body:
        for paragraph in _text(body).splitlines():
            if paragraph.strip():
                story.append(
                    Paragraph(
                        _pdf_markup_text(paragraph, base_font=base_font, cjk_font=cjk_font),
                        styles["BodyText"],
                    )
                )
    else:
        story.append(Paragraph("No report body was provided.", styles["BodyText"]))

    doc.build(story)
    return _published_response(
        payload=output.getvalue(),
        name=_ensure_name(name, default="generated.pdf", suffix=".pdf"),
        mime=_PDF_MIME,
        source="create_pdf_report",
    )
